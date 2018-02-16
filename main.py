from pptx import Presentation
from pptx.util import Inches, Pt
import easygui

def extract_text(input_file):
    prs = Presentation(input_file)
    text_runs = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    slide_text.append(run.text)
        text_runs.append(slide_text)
    final_text = []
    slide_text=""
    for slide in text_runs:
        for line in slide:
            slide_text+=line
            slide_text+="\n"
        final_text.append(slide_text)
        slide_text=""
    return final_text

def create_presentation(list_of_slides_text,output_file):
    prs = Presentation()
    for text_slide in list_of_slides_text:
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        left = width = height = Inches(0.5)
        top = Inches(0.1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p= tf.add_paragraph()
        p.text = text_slide
        p.font.size = Pt(32)
    prs.save(output_file)

if __name__ == "__main__":
    extractedtxt = []
    while True:
        choice = easygui.buttonbox('Choose "Add Slide" to add song slides to the powerpoint, and press export when done.', title='eazymworship v0.1', choices=["Add Slide","Export","Quit"])
        if choice == "Add Slide":
            input_filename = easygui.fileopenbox('Choose a Powerpoint File')
            extractedtxt+= extract_text(input_filename)
        elif choice == "Export":
            output_filename = easygui.enterbox("Save as... (without .pptx)")
            create_presentation(extractedtxt, output_filename+".pptx")
            break
        elif choice == "Quit":
            break
        else:
            break
